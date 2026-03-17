#!/usr/bin/env python3
"""
build_pptx.py — Build a PPTX from PDF page images + LLM-extracted text layout.

Usage:
  python3 build_pptx.py <input.json> <output.pptx>

input.json schema:
{
  "pdf_path": "/tmp/input.pdf",
  "slides": [
    {
      "page_num": 1,
      "elements": [
        {
          "type": "text",
          "content": "Hello World",
          "x": 10.5,   // % of slide width
          "y": 20.0,   // % of slide height
          "w": 60.0,   // % of slide width
          "h": 8.0,    // % of slide height
          "fontSize": 24,
          "bold": false,
          "color": "#FFFFFF",
          "align": "left"
        }
      ]
    }
  ]
}
"""

import sys
import json
import os
import tempfile
import traceback

def hex_to_rgb(hex_color):
    """Convert #RRGGBB to (R, G, B) tuple."""
    h = hex_color.lstrip('#')
    if len(h) == 3:
        h = ''.join(c*2 for c in h)
    if len(h) != 6:
        return (0, 0, 0)
    try:
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
    except ValueError:
        return (0, 0, 0)

def build_pptx(input_json_path, output_pptx_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pdf2image import convert_from_path
    from PIL import Image
    import io

    with open(input_json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    pdf_path = data['pdf_path']
    slides_data = data['slides']

    # Slide dimensions: 16:9 widescreen (13.33" x 7.5")
    SLIDE_W_IN = 13.33
    SLIDE_H_IN = 7.5
    SLIDE_W_EMU = int(SLIDE_W_IN * 914400)
    SLIDE_H_EMU = int(SLIDE_H_IN * 914400)

    # Render PDF pages to images at high DPI
    print(f"Rendering {len(slides_data)} PDF page(s) to images...", flush=True)
    try:
        page_nums = [s['page_num'] for s in slides_data]
        first_page = min(page_nums)
        last_page = max(page_nums)
        images = convert_from_path(
            pdf_path,
            dpi=150,
            first_page=first_page,
            last_page=last_page,
            fmt='png',
            thread_count=2,
        )
        # Map page_num -> PIL image
        page_images = {}
        for i, img in enumerate(images):
            page_images[first_page + i] = img
        print(f"Rendered {len(page_images)} page image(s)", flush=True)
    except Exception as e:
        print(f"WARNING: PDF rendering failed: {e}", flush=True)
        page_images = {}

    # Create presentation
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)

    # Use blank slide layout (index 6 is typically blank)
    blank_layout = prs.slide_layouts[6]

    ALIGN_MAP = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY,
    }

    for slide_data in slides_data:
        page_num = slide_data['page_num']
        elements = slide_data.get('elements', [])

        slide = prs.slides.add_slide(blank_layout)

        # ── Add page image as full-slide background ──────────────────────
        if page_num in page_images:
            img = page_images[page_num]
            # Save to bytes
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            # Add as picture covering the full slide
            pic = slide.shapes.add_picture(
                img_bytes,
                left=Emu(0),
                top=Emu(0),
                width=Emu(SLIDE_W_EMU),
                height=Emu(SLIDE_H_EMU),
            )
            # Move picture to back (z-order 0)
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            print(f"  Slide {page_num}: background image added ({img.width}x{img.height}px)", flush=True)
        else:
            print(f"  Slide {page_num}: no background image (white)", flush=True)

        # ── Overlay text elements ────────────────────────────────────────
        for elem in elements:
            if elem.get('type') != 'text':
                continue
            content = (elem.get('content') or '').strip()
            if not content:
                continue

            # Convert percentage positions to EMU
            x_pct = float(elem.get('x', 0))
            y_pct = float(elem.get('y', 0))
            w_pct = float(elem.get('w', 80))
            h_pct = float(elem.get('h', 5))

            left   = Emu(int(x_pct / 100 * SLIDE_W_EMU))
            top    = Emu(int(y_pct / 100 * SLIDE_H_EMU))
            width  = Emu(max(int(w_pct / 100 * SLIDE_W_EMU), 914400))  # min 1"
            height = Emu(max(int(h_pct / 100 * SLIDE_H_EMU), 182880))  # min 0.2"

            # Clamp to slide bounds
            if left + width > Emu(SLIDE_W_EMU):
                width = Emu(SLIDE_W_EMU) - left
            if top + height > Emu(SLIDE_H_EMU):
                height = Emu(SLIDE_H_EMU) - top

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.alignment = ALIGN_MAP.get(elem.get('align', 'left'), PP_ALIGN.LEFT)

            run = p.add_run()
            run.text = content

            # Font properties
            font = run.font
            raw_size = elem.get('fontSize', 14)
            try:
                font_size = max(8, min(int(float(raw_size)), 96))
            except (TypeError, ValueError):
                font_size = 14
            font.size = Pt(font_size)
            font.bold = bool(elem.get('bold', False))

            raw_color = elem.get('color', '#000000') or '#000000'
            r, g, b = hex_to_rgb(raw_color)
            font.color.rgb = RGBColor(r, g, b)

            # Transparent text box background
            from pptx.util import Pt as _Pt
            from lxml import etree
            spPr = txBox._element.spPr
            # Remove any fill so it's transparent
            noFill = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')

    prs.save(output_pptx_path)
    print(f"Saved PPTX: {output_pptx_path}", flush=True)

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: build_pptx.py <input.json> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    try:
        build_pptx(sys.argv[1], sys.argv[2])
    except Exception as e:
        traceback.print_exc()
        sys.exit(1)
